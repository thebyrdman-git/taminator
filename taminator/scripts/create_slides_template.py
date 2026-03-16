#!/usr/bin/env python3
"""
Create a Taminator graphical template for Google Slides.
Output: TAMINATOR_DEMO_TEMPLATE.pptx
Upload to Google Slides: File → Open → Upload, or drag the file to drive.google.com and open with Google Slides.
Requires: pip install python-pptx
"""
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    raise

# Taminator / Red Hat inspired colors (RGB 0-255)
REDHAT_RED = RGBColor(0xCC, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG = RGBColor(0xFA, 0xFA, 0xFA)
TITLE_BAR_BG = RGBColor(0x1a, 0x1a, 0x2e)  # dark blue-gray

OUTPUT_DIR = Path(__file__).resolve().parent.parent
OUTPUT_FILE = OUTPUT_DIR / "TAMINATOR_DEMO_TEMPLATE.pptx"


def set_shape_fill(shape, rgb):
    if shape.fill.type is None:
        shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def add_title_slide(prs):
    """Slide 1: Title slide with Taminator branding."""
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    # Dark bar at top
    top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))  # rectangle
    top.fill.solid()
    top.fill.fore_color.rgb = TITLE_BAR_BG
    top.line.fill.background()
    tf = top.text_frame
    tf.paragraphs[0].text = "Taminator"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RFE & Bug Report Generator"
    p.font.size = Pt(44)
    p.font.color.rgb = DARK_GRAY
    p.font.bold = True
    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.8))
    sub.text_frame.paragraphs[0].text = "for Red Hat TAMs"
    sub.text_frame.paragraphs[0].font.size = Pt(24)
    sub.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    # Footer line
    foot = slide.shapes.add_shape(1, Inches(0.5), Inches(6.8), Inches(9), Inches(0.02))
    foot.fill.solid()
    foot.fill.fore_color.rgb = REDHAT_RED
    foot.line.fill.background()
    # Footer text
    foot_t = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
    foot_t.text_frame.paragraphs[0].text = "Demo — [Your name] — [Date]"
    foot_t.text_frame.paragraphs[0].font.size = Pt(12)
    foot_t.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY


def add_content_slide(prs, title_text, bullet_points=None):
    """Add a content slide with title bar and optional bullets."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BAR_BG
    bar.line.fill.background()
    tb = bar.text_frame
    tb.paragraphs[0].text = title_text
    tb.paragraphs[0].font.size = Pt(24)
    tb.paragraphs[0].font.color.rgb = WHITE
    tb.paragraphs[0].font.bold = True
    # Content
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    if bullet_points:
        for i, line in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line.strip().replace("|", "\n")
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_content_slide(prs, "What is Taminator?", [
        "Generates and maintains RFE and bug reports for Red Hat TAMs",
        "Uses JIRA and case data → consistent markdown reports",
        "Can post to customer portal groups",
        "Saves 2–3 hours per customer per week vs manual tracking",
    ])
    add_content_slide(prs, "How to get it", [
        "Desktop app: AppImage (Linux) or DMG (macOS) from GitLab releases",
        "From repo: git clone for CLI or development",
        "Official: gitlab.cee.redhat.com/jbyrd/taminator (Red Hat VPN)",
    ])
    add_content_slide(prs, "Desktop app — no terminal required", [
        "Double-click the app; it opens in its own window",
        "Full UI: Report Manager • Check/Update • Library • Settings • User Guide",
        "Linux: AppImage | macOS: DMG (Intel + Apple Silicon) | Windows: installer",
    ])
    add_content_slide(prs, "Core workflow", [
        "1. Create a report (Report Manager — add account, SBR Group, account numbers)",
        "2. Check report — compare to JIRA (no file changes)",
        "3. Update report — write current JIRA statuses into file (backup first)",
        "4. Post (optional) — post to customer portal group",
    ])
    add_content_slide(prs, "JIRA: Red Hat + Cloud", [
        "Red Hat JIRA (default): issues.redhat.com — Personal Access Token",
        "JIRA Cloud: your-tenant.atlassian.net — Email + API token",
        "Set in Settings → Token configuration → JIRA instance",
        "Same Check/Update flow for both",
    ])
    add_content_slide(prs, "Security & credentials", [
        "Tokens in ~/.config/taminator/ui_tokens.json (encoded)",
        "Optional: HashiCorp Vault — VAULT_ADDR + VAULT_TOKEN",
        "VPN required for Red Hat JIRA (skipped for JIRA Cloud)",
    ])
    add_content_slide(prs, "Demo outline", [
        "1. Open Taminator (desktop app)",
        "2. Report Manager — report creation / accounts",
        "3. Check/Update Reports — run Check or Update",
        "4. Library — browse reports | 5. Settings — VPN, tokens | 6. User Guide",
    ])
    add_content_slide(prs, "Links & feedback", [
        "Releases: gitlab.cee.redhat.com/jbyrd/taminator/-/releases",
        "User guide: repo README",
        "Issues: gitlab.cee.redhat.com/jbyrd/taminator/-/issues",
    ])
    add_content_slide(prs, "Questions?", [
        "Taminator — RFE and Bug Report Generator for Red Hat TAMs",
    ])

    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")
    print("Upload to Google Slides: File → Open → Upload (or drag to Google Drive and open with Google Slides).")


if __name__ == "__main__":
    main()
